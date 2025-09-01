#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_enhanced_presentation():
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Définir les couleurs du thème professionnel
    primary_color = RGBColor(59, 130, 246)      # Bleu principal
    secondary_color = RGBColor(16, 185, 129)    # Vert succès
    accent_color = RGBColor(245, 158, 11)       # Orange accent
    warning_color = RGBColor(239, 68, 68)       # Rouge warning
    dark_color = RGBColor(31, 41, 55)           # Gris foncé
    light_color = RGBColor(249, 250, 251)       # Gris clair
    
    # SLIDE 1 : Page de Titre Améliorée
    slide_layout = prs.slide_layouts[6]  # Layout vide pour plus de contrôle
    slide = prs.slides.add_slide(slide_layout)
    
    # Ajouter un rectangle de fond
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = light_color
    background.line.fill.background()
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "Maintenance et Évolution"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Plateforme de Cours en Ligne\nBloc 4 - Gestion du Cycle de Vie des Applications"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = dark_color
    
    # Ajouter des formes décoratives
    # Cercle décoratif en haut à gauche
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), Inches(0.5),
        Inches(1), Inches(1)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = secondary_color
    circle1.line.fill.background()
    
    # Cercle décoratif en bas à droite
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(11.83), Inches(6),
        Inches(1), Inches(1)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = accent_color
    circle2.line.fill.background()
    
    # SLIDE 2 : Vue d'Ensemble avec Style Amélioré
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre avec style
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Vue d'Ensemble du Projet"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Contenu en colonnes
    # Colonne gauche
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(5))
    left_frame = left_content.text_frame
    left_frame.text = """🎯 Contexte et Objectifs

• Problématique : Créer une plateforme moderne d'apprentissage en ligne
• Solution : Architecture microservices avec technologies modernes
• Public cible : Étudiants, instructeurs, administrateurs

📊 Chiffres Clés

• 7 microservices interconnectés
• 15,000+ utilisateurs actifs
• 500+ cours disponibles
• 99.95% uptime en production"""
    
    # Colonne droite
    right_content = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(5))
    right_frame = right_content.text_frame
    right_frame.text = """🏗️ Architecture Technique

Frontend :
• React + TypeScript
• Material-UI
• Responsive Design

Backend :
• Node.js + Express
• MongoDB + Mongoose
• JWT Authentication

Monitoring :
• Prometheus + Grafana
• AlertManager
• Health Checks

DevOps :
• Docker + Docker Compose
• GitHub Actions CI/CD
• SonarQube Quality"""
    
    # Ajouter des icônes décoratives
    icon1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(6), Inches(0.1))
    icon1.fill.solid()
    icon1.fill.fore_color.rgb = secondary_color
    icon1.line.fill.background()
    
    icon2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(6), Inches(0.1))
    icon2.fill.solid()
    icon2.fill.fore_color.rgb = accent_color
    icon2.line.fill.background()
    
    # SLIDE 3 : C4.1.1 avec Style Amélioré
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.1.1 - Processus de Mise à Jour des Dépendances"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative sous le titre
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = secondary_color
    line.line.fill.background()
    
    # Contenu en sections
    # Section 1
    section1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(2.5))
    section1_frame = section1.text_frame
    section1_frame.text = """🔄 Fréquence des Mises à Jour

• Processus automatisé hebdomadaire avec Dependabot
• Vérification manuelle mensuelle pour les dépendances critiques
• Mises à jour trimestrielles pour les composants majeurs

🛠️ Type de Mise à Jour

• Automatique : Dependabot + GitHub Actions
• Manuel : Dépendances critiques et sécurité
• Validation : Tests automatisés avant déploiement"""
    
    # Section 2 - Tableau stylisé
    section2 = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(4.5))
    section2_frame = section2.text_frame
    section2_frame.text = """🎯 Périmètre Concerné

┌─────────────────┬─────────────┬─────────────┬─────────────┐
│   Composant     │  Fréquence  │    Type     │ Responsable │
├─────────────────┼─────────────┼─────────────┼─────────────┤
│   Frontend      │ Hebdomadaire│ Automatique │ Dependabot  │
│ Backend Services│ Hebdomadaire│ Automatique │ Dependabot  │
│ Base de données │Trimestrielle│   Manuel    │   DevOps    │
│   Monitoring    │  Mensuelle  │ Automatique │ Dependabot  │
└─────────────────┴─────────────┴─────────────┴─────────────┘"""
    
    # SLIDE 4 : C4.1.2 avec Métriques Visuelles
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.1.2 - Système de Supervision et d'Alerte"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    # Contenu en 3 colonnes
    # Colonne 1 - Système
    col1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4.5))
    col1_frame = col1.text_frame
    col1_frame.text = """🎯 Système Adapté

Prometheus + Grafana + AlertManager
pour microservices

• Collecte de métriques en temps réel
• Visualisation via dashboards personnalisés
• Alertes intelligentes avec seuils configurables
• Health checks toutes les 30 secondes"""
    
    # Colonne 2 - Métriques
    col2 = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4.5))
    col2_frame = col2.text_frame
    col2_frame.text = """🔍 Sondes Mises en Place

Métriques Système :
• CPU, mémoire, disque, réseau
• Temps de réponse des APIs
• Taux d'erreur par service

Métriques Métier :
• Utilisateurs actifs
• Cours consultés
• Transactions de paiement"""
    
    # Colonne 3 - Critères
    col3 = slide.shapes.add_textbox(Inches(9.5), Inches(2), Inches(3.33), Inches(4.5))
    col3_frame = col3.text_frame
    col3_frame.text = """📈 Critères de Qualité

• Temps de réponse < 200ms
• Uptime > 99.9%
• Taux d'erreur < 1%
• Utilisation CPU < 80%

✅ Surveillance Disponibilité

• Alertes automatiques en cas de panne
• Escalade vers l'équipe technique
• Monitoring 24/7"""
    
    # SLIDE 5 : C4.2.1 avec Fiche d'Anomalie Stylisée
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.2.1 - Processus de Collecte et Consignation des Anomalies"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = warning_color
    line.line.fill.background()
    
    # Processus
    process_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(2))
    process_frame = process_box.text_frame
    process_frame.text = """📋 Processus Structuré et Adapté

• GitHub Issues comme système centralisé
• Template standardisé pour tous les bugs
• Workflow automatisé de suivi
• Intégration avec le pipeline CI/CD"""
    
    # Fiche d'anomalie stylisée
    anomaly_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(4.5))
    anomaly_frame = anomaly_box.text_frame
    anomaly_frame.text = """🐛 Fiche d'Anomalie Complète

Bug #123 - Affichage des cours utilisateur

📅 Date détection : 2025-01-10
🚨 Sévérité : Critique
👥 Impact : 100% des utilisateurs

📝 Étapes de reproduction :
1. Se connecter avec un compte utilisateur
2. Aller sur le profil utilisateur
3. Cliquer sur l'onglet "Mes Cours"
4. Voir le message "Vous n'êtes inscrit à aucun cours"

🔍 Analyse et Recommandations
• Cause identifiée : Modèle UserCourse non initialisé
• Correction prévue : Validation et fallback
• Tests post-correction : 4 tests unitaires ajoutés"""
    
    # SLIDE 6 : C4.2.2 avec Code Stylisé
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.2.2 - Traitement d'une Anomalie Détectée"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = secondary_color
    line.line.fill.background()
    
    # Bénéfice CI/CD
    benefit_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(1.5))
    benefit_frame = benefit_box.text_frame
    benefit_frame.text = """🚀 Bénéfice du Processus CI/CD

• Intégration continue pour la correction
• Branche de correctif créée automatiquement
• Tests automatisés avant merge
• Déploiement automatique après validation"""
    
    # Code avant/après
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(3))
    code_frame = code_box.text_frame
    code_frame.text = """📝 Description de l'Action Corrective

AVANT (code problématique) :
const getUserCourses = async (userId) => {
  const userCourses = await UserCourse.find({ userId });
  return userCourses.courses; // ❌ Erreur si null
};

APRÈS (code corrigé) :
const getUserCourses = async (userId) => {
  try {
    const userCourses = await UserCourse.find({ userId });
    if (!userCourses || userCourses.length === 0) {
      return [];
    }
    return userCourses.map(uc => uc.courses).flat();
  } catch (error) {
    logger.error('Erreur lors de la récupération des cours:', error);
    return [];
  }
};

✅ Résolution de l'Anomalie
• Code corrigé et testé
• Déploiement automatique en production
• Monitoring renforcé pour prévenir la récurrence"""
    
    # SLIDE 7 : C4.3.1 avec Recommandations Visuelles
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.3.1 - Recommandations d'Amélioration Raisonnées"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    # Recommandations en 3 colonnes
    # Recommandation 1
    rec1 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(4.5))
    rec1_frame = rec1.text_frame
    rec1_frame.text = """💾 Cache Redis pour les Performances

Gain attendu : 70% réduction temps de réponse
Coût d'implémentation : 2 semaines
ROI estimé : 40% économie sur les coûts serveur

✅ Réalisme et Faisabilité
• Technologies maîtrisées par l'équipe
• Ressources disponibles
• Impact mesurable"""
    
    # Recommandation 2
    rec2 = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(4.5))
    rec2_frame = rec2.text_frame
    rec2_frame.text = """📱 Interface Progressive Web App (PWA)

Gain attendu : 25% augmentation du temps passé
Coût d'implémentation : 3 semaines
ROI estimé : Amélioration de l'engagement utilisateur

🎯 Renforcement de l'Attractivité
• Expérience utilisateur améliorée
• Sécurité renforcée
• Performance optimisée"""
    
    # Recommandation 3
    rec3 = slide.shapes.add_textbox(Inches(9.5), Inches(2), Inches(3.33), Inches(4.5))
    rec3_frame = rec3.text_frame
    rec3_frame.text = """🔐 Authentification Multi-Facteurs (2FA)

Gain attendu : 90% réduction des tentatives de hack
Coût d'implémentation : 1 semaine
ROI estimé : Conformité RGPD renforcée

🛡️ Sécurité Renforcée
• Protection contre les attaques
• Conformité réglementaire
• Confiance utilisateur"""
    
    # SLIDE 8 : C4.3.2 avec Journal des Versions Stylisé
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.3.2 - Exemple de Journal des Versions"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = secondary_color
    line.line.fill.background()
    
    # Version header
    version_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(0.8))
    version_frame = version_box.text_frame
    version_frame.text = "📋 Version 2.1.0 - 2025-01-15"
    version_para = version_frame.paragraphs[0]
    version_para.font.size = Pt(20)
    version_para.font.bold = True
    version_para.font.color.rgb = accent_color
    
    # Contenu en 3 sections
    # Nouvelles fonctionnalités
    features_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(3.5))
    features_frame = features_box.text_frame
    features_frame.text = """🚀 Nouvelles Fonctionnalités

• Système de recommandation basé sur l'IA
• Notifications push temps réel
• Interface d'administration avancée
• Export des données en CSV"""
    
    # Anomalies corrigées
    bugs_box = slide.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(3.5))
    bugs_frame = bugs_box.text_frame
    bugs_frame.text = """🔧 Anomalies Corrigées

• #123 - Affichage des cours utilisateur
• #124 - Upload vidéos > 100MB
• #125 - Synchronisation paiements Stripe"""
    
    # Améliorations
    improvements_box = slide.shapes.add_textbox(Inches(9.5), Inches(3), Inches(3.33), Inches(3.5))
    improvements_frame = improvements_box.text_frame
    improvements_frame.text = """⚡ Améliorations

• Performance : +40% temps de chargement
• Sécurité : Mise à jour dépendances critiques
• UX : Design responsive amélioré

✅ Actions Correctives Documentées
• Code source des corrections
• Tests unitaires ajoutés
• Documentation mise à jour
• Monitoring renforcé"""
    
    # SLIDE 9 : C4.3.3 avec Collaboration Stylisée
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "C4.3.3 - Exemple de Problème Résolu en Collaboration"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = warning_color
    line.line.fill.background()
    
    # Contexte client
    context_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(1.5))
    context_frame = context_box.text_frame
    context_frame.text = """📋 Contexte du Retour Client

• Problème : "Impossible de finaliser l'achat d'un cours"
• Client : Marie D. (utilisatrice premium)
• Impact : Perte de revenus, frustration client
• Ticket : #SUPPORT-2025-001"""
    
    # Diagnostic technique
    diagnostic_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(6), Inches(1.5))
    diagnostic_frame = diagnostic_box.text_frame
    diagnostic_frame.text = """🔍 Diagnostic Technique

2025-01-15 14:25:12 ERROR [payment-service] Stripe API Error
Error: Invalid API key provided
    at Stripe.createPaymentIntent (/app/src/services/stripe.service.js:45)

Cause identifiée : Variable d'environnement STRIPE_SECRET_KEY corrompue"""
    
    # Solution et contribution
    solution_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(2.5))
    solution_frame = solution_box.text_frame
    solution_frame.text = """🛠️ Solution Fournie

Actions correctives :
1. Régénération de la clé Stripe
2. Mise à jour de la variable d'environnement
3. Redémarrage du service de paiement
4. Test de validation de l'endpoint

👥 Contribution des Différents Acteurs

Support (Niveau 1) : Réception et qualification du problème
Équipe Technique : Diagnostic et correction technique
Client : Test de validation post-correction, Feedback positif (5/5 ⭐)"""
    
    # SLIDE 10 : Métriques avec Graphiques Visuels
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Métriques de Maintenance et Évolution"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = secondary_color
    line.line.fill.background()
    
    # Tableau de métriques stylisé
    metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(2.5))
    metrics_frame = metrics_box.text_frame
    metrics_frame.text = """🎯 Indicateurs de Performance

┌─────────────────────────┬─────────────┬─────────────┬─────────┐
│        Métrique         │   Objectif  │   Actuel    │  Statut │
├─────────────────────────┼─────────────┼─────────────┼─────────┤
│ Temps de résolution bugs│    < 4h     │   2h30      │   ✅    │
│ Disponibilité système   │   > 99.9%   │   99.95%    │   ✅    │
│ Taux de satisfaction    │   > 4.5/5   │    4.8/5    │   ✅    │
│ Fréquence des mises à jour│ Hebdomadaire│  Respecté   │   ✅    │
│ Couverture monitoring   │    100%     │    100%     │   ✅    │
└─────────────────────────┴─────────────┴─────────────┴─────────┘"""
    
    # Points forts
    strengths_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.33), Inches(2))
    strengths_frame = strengths_box.text_frame
    strengths_frame.text = """🏆 Points Forts de la Maintenance

🔄 Processus Automatisé : CI/CD pour les déploiements, Monitoring en temps réel, Alertes intelligentes
📈 Qualité Continue : Tests automatisés (88% couverture), Analyse statique avec SonarQube, Audit de sécurité régulier
👥 Collaboration Efficace : Support client réactif, Documentation à jour, Formation de l'équipe"""
    
    # SLIDE 11 : Conclusion avec Checklist Visuelle
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion et Perspectives"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()
    
    # Compétences maîtrisées
    skills_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(3.5))
    skills_frame = skills_box.text_frame
    skills_frame.text = """✅ Compétences Maîtrisées

🚀 Bloc 4 - Maintenance et Évolution
• ✅ C4.1.1 - Processus de mise à jour des dépendances
• ✅ C4.1.2 - Système de supervision et d'alerte
• ✅ C4.2.1 - Collecte et consignation des anomalies
• ✅ C4.2.2 - Traitement d'anomalies avec CI/CD
• ✅ C4.3.1 - Recommandations d'amélioration
• ✅ C4.3.2 - Journal des versions
• ✅ C4.3.3 - Collaboration avec le support client"""
    
    # Évolution et apprentissages
    evolution_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.83), Inches(3.5))
    evolution_frame = evolution_box.text_frame
    evolution_frame.text = """🔮 Évolution Continue

📈 Améliorations Planifiées
• Machine Learning pour la détection d'anomalies
• Auto-scaling basé sur la charge
• Monitoring prédictif avec IA

💡 Apprentissages Clés
• Maintenance proactive vs réactive
• Importance du monitoring en temps réel
• Collaboration entre équipes
• Documentation comme investissement"""
    
    # SLIDE 12 : Questions et Réponses avec Style Final
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Questions et Réponses"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Ligne décorative
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.33), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = secondary_color
    line.line.fill.background()
    
    # Points de discussion
    discussion_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(3.5))
    discussion_frame = discussion_box.text_frame
    discussion_frame.text = """🎯 Points de Discussion

🔧 Maintenance et Monitoring
• Stratégie de mise à jour des dépendances
• Configuration du système de supervision
• Gestion des alertes et escalade

🐛 Gestion des Anomalies
• Processus de détection et qualification
• Intégration avec le CI/CD
• Communication avec les clients

🚀 Évolution et Amélioration
• Critères de sélection des améliorations
• Mesure du ROI des évolutions
• Planification des releases"""
    
    # Contact et remerciements
    contact_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.83), Inches(3.5))
    contact_frame = contact_box.text_frame
    contact_frame.text = """📞 Contact et Ressources

• Repository GitHub : [URL]
• Documentation technique : [URL]
• Dashboard monitoring : [URL]
• Support client : [Email]

🎉 Merci pour votre attention !

Prêt pour les questions !"""
    
    # Ajouter des formes décoratives finales
    # Cercle décoratif en bas à gauche
    circle3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.5), Inches(6),
        Inches(0.8), Inches(0.8)
    )
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = secondary_color
    circle3.line.fill.background()
    
    # Cercle décoratif en bas à droite
    circle4 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(12), Inches(6),
        Inches(0.8), Inches(0.8)
    )
    circle4.fill.solid()
    circle4.fill.fore_color.rgb = accent_color
    circle4.line.fill.background()
    
    # Sauvegarder la présentation
    filename = "Presentation_Bloc4_Enhanced.pptx"
    prs.save(filename)
    print(f"✅ Présentation améliorée créée avec succès : {filename}")
    print(f"📁 Fichier sauvegardé dans : {os.path.abspath(filename)}")
    
    return filename

if __name__ == "__main__":
    try:
        filename = create_enhanced_presentation()
        print(f"\n🎉 Présentation PowerPoint améliorée générée avec succès !")
        print(f"📊 Contenu : 12 slides avec style professionnel")
        print(f"🎨 Améliorations : Couleurs cohérentes, formes décoratives, mise en page optimisée")
        print(f"🎯 Prêt pour votre présentation de certification")
    except Exception as e:
        print(f"❌ Erreur lors de la création de la présentation : {e}")
